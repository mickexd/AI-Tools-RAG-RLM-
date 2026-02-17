#!/usr/bin/env python3
"""
Document RAG MCP Server for Claude
Handles PDF, DOCX, XLSX, PPTX, TXT, MD and more with hybrid search
CPU-OPTIMIZED: Tuned for maximum CPU performance
CROSS-PLATFORM: Works on Windows, macOS, and Linux
"""

import asyncio
import json
import sys
import os
import hashlib
import platform
from datetime import datetime
from pathlib import Path
from typing import Any, List, Optional, Dict, Tuple
from dataclasses import dataclass
from collections import defaultdict

# Redirect stdout to stderr for debugging (MCP uses stdout for JSON-RPC)
import builtins

original_print = builtins.print


def debug_print(*args, **kwargs):
    kwargs["file"] = sys.stderr
    original_print(*args, **kwargs)


builtins.print = debug_print

# CPU optimization environment variables
cpu_cores = os.cpu_count() or 4
os.environ["TORCH_NUM_THREADS"] = str(cpu_cores)
os.environ["OMP_NUM_THREADS"] = str(cpu_cores)
os.environ["OPENBLAS_NUM_THREADS"] = str(cpu_cores)
os.environ["MKL_NUM_THREADS"] = str(cpu_cores)
os.environ["VECLIB_MAXIMUM_THREADS"] = str(cpu_cores)
os.environ["NUMEXPR_NUM_THREADS"] = str(cpu_cores)

try:
    # MCP imports
    from mcp.server import Server
    from mcp.server.stdio import stdio_server
    from mcp.types import Tool, TextContent

    # Document parsing libraries
    UNSTRUCTURED_AVAILABLE = False
    PYPDF2_AVAILABLE = False
    PYTHON_DOCX_AVAILABLE = False
    OPENPYXL_AVAILABLE = False
    PYTHON_PPTX_AVAILABLE = False
    PANDAS_AVAILABLE = False
    RANK_BM25_AVAILABLE = False

    try:
        from unstructured.partition.pdf import partition_pdf
        from unstructured.partition.docx import partition_docx
        from unstructured.partition.auto import partition

        UNSTRUCTURED_AVAILABLE = True
    except ImportError:
        pass

    try:
        import PyPDF2

        PYPDF2_AVAILABLE = True
    except ImportError:
        pass

    try:
        import docx

        PYTHON_DOCX_AVAILABLE = True
    except ImportError:
        pass

    try:
        import openpyxl

        OPENPYXL_AVAILABLE = True
    except ImportError:
        pass

    try:
        from pptx import Presentation

        PYTHON_PPTX_AVAILABLE = True
    except ImportError:
        pass

    try:
        import pandas as pd

        PANDAS_AVAILABLE = True
    except ImportError:
        pass

    try:
        from rank_bm25 import BM25Okapi

        RANK_BM25_AVAILABLE = True
    except ImportError:
        pass

    # LanceDB and embeddings
    import lancedb
    from sentence_transformers import SentenceTransformer
    import torch
    import pyarrow as pa

    # Text chunking
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    # Tokenization for BM25
    import nltk

    try:
        nltk.data.find("tokenizers/punkt")
    except LookupError:
        nltk.download("punkt", quiet=True)

    print("✓ Imports successful", flush=True)
except Exception as e:
    print(f"✗ Import error: {e}", flush=True)
    sys.exit(1)


# Cross-platform configuration
def get_documents_dir():
    """Get cross-platform documents directory"""
    system = platform.system()
    if system == "Windows":
        return Path.home() / ".claude_documents"
    elif system == "Darwin":  # macOS
        return Path.home() / "Library" / "Application Support" / "ClaudeDocuments"
    else:  # Linux
        return Path.home() / ".local" / "share" / "claude_documents"


# Configuration
DOCUMENTS_DIR = get_documents_dir()
DOCUMENTS_DIR.mkdir(parents=True, exist_ok=True)
DB_PATH = str(DOCUMENTS_DIR / "lancedb")
UPLOADS_DIR = DOCUMENTS_DIR / "uploads"
UPLOADS_DIR.mkdir(exist_ok=True)
VERSIONS_DIR = DOCUMENTS_DIR / "versions"
VERSIONS_DIR.mkdir(exist_ok=True)
BM25_INDEX_FILE = DOCUMENTS_DIR / "bm25_index.json"

# Supported file formats
SUPPORTED_FORMATS = {
    # Documents
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "doc",
    # Spreadsheets
    ".xlsx": "excel",
    ".xls": "excel",
    ".csv": "csv",
    # Presentations
    ".pptx": "powerpoint",
    ".ppt": "powerpoint",
    # Text files
    ".txt": "text",
    ".md": "markdown",
    ".markdown": "markdown",
    ".rst": "text",
    ".html": "html",
    ".htm": "html",
    ".xml": "xml",
    ".json": "json",
    ".yaml": "yaml",
    ".yml": "yaml",
    ".log": "text",
    ".rtf": "text",
    ".tex": "text",
    ".css": "text",
    ".js": "text",
    ".py": "text",
    ".java": "text",
    ".cpp": "text",
    ".c": "text",
    ".h": "text",
    ".sql": "text",
    ".sh": "text",
    ".bat": "text",
    ".ps1": "text",
}

# Initialize server
server = Server("document_rag")

# Global state
db = None
documents_table = None
versions_table = None
embed_model = None
text_splitter = None
bm25_index = None
bm25_corpus = []
bm25_doc_ids = []


def init_database():
    """Initialize LanceDB and embedding model - CPU optimized"""
    global db, documents_table, versions_table, embed_model, text_splitter

    try:
        # Force CPU for stability
        device = "cpu"
        cpu_count = os.cpu_count() or 4

        print(f"═" * 60, flush=True)
        print(f"DOCUMENT RAG SERVER - {platform.system()}", flush=True)
        print(f"═" * 60, flush=True)
        print(f"Device: CPU (optimized for performance)", flush=True)
        print(f"Available CPU cores: {cpu_count}", flush=True)
        print(f"Storage: {DOCUMENTS_DIR}", flush=True)

        # Load embedding model
        print(f"\nInitializing embedding model...", flush=True)
        embed_model = SentenceTransformer(
            "all-MiniLM-L6-v2",
            device=device,
            model_kwargs={"torch_dtype": torch.float32},
        )
        embed_model = embed_model.to(device)
        embed_model.eval()

        print(f"✓ Embedding model loaded", flush=True)
        print(f"  Embedding dimension: 384", flush=True)

        # Initialize text splitter
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""],
        )
        print(f"✓ Text splitter initialized (chunk_size=1000, overlap=200)", flush=True)

        # Connect to LanceDB
        db = lancedb.connect(DB_PATH)
        print(f"✓ Connected to LanceDB at {DB_PATH}", flush=True)

        # Create or open documents table
        try:
            documents_table = db.open_table("documents")
            row_count = len(documents_table)
            print(f"✓ Opened existing documents table ({row_count} chunks)", flush=True)
        except Exception:
            schema = pa.schema(
                [
                    pa.field("id", pa.string()),
                    pa.field("text", pa.string()),
                    pa.field("vector", pa.list_(pa.float32(), 384)),
                    pa.field("document_id", pa.string()),
                    pa.field("document_name", pa.string()),
                    pa.field("chunk_index", pa.int32()),
                    pa.field("page_number", pa.int32()),
                    pa.field("total_pages", pa.int32()),
                    pa.field("timestamp", pa.string()),
                    pa.field("file_type", pa.string()),
                    pa.field("version", pa.int32()),
                    pa.field("keywords", pa.string()),  # For BM25
                    pa.field("metadata", pa.string()),
                ]
            )
            documents_table = db.create_table("documents", schema=schema)
            print(f"✓ Created new documents table", flush=True)

        # Create or open versions table
        try:
            versions_table = db.open_table("document_versions")
            print(f"✓ Opened existing versions table", flush=True)
        except Exception:
            version_schema = pa.schema(
                [
                    pa.field("version_id", pa.string()),
                    pa.field("document_id", pa.string()),
                    pa.field("document_name", pa.string()),
                    pa.field("version_number", pa.int32()),
                    pa.field("timestamp", pa.string()),
                    pa.field("file_hash", pa.string()),
                    pa.field("file_path", pa.string()),
                    pa.field("chunk_count", pa.int32()),
                    pa.field("metadata", pa.string()),
                ]
            )
            versions_table = db.create_table("document_versions", schema=version_schema)
            print(f"✓ Created new versions table", flush=True)

        # Initialize BM25 index
        load_bm25_index()

        print(f"═" * 60, flush=True)

    except Exception as e:
        print(f"✗ Database initialization error: {e}", flush=True)
        import traceback

        traceback.print_exc(file=sys.stderr)
        raise


def tokenize_text(text: str) -> List[str]:
    """Tokenize text for BM25"""
    try:
        from nltk.tokenize import word_tokenize

        return word_tokenize(text.lower())
    except:
        # Fallback simple tokenization
        return text.lower().split()


def load_bm25_index():
    """Load or rebuild BM25 index"""
    global bm25_index, bm25_corpus, bm25_doc_ids

    try:
        if BM25_INDEX_FILE.exists():
            with open(BM25_INDEX_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                bm25_corpus = data.get("corpus", [])
                bm25_doc_ids = data.get("doc_ids", [])
                if bm25_corpus:
                    bm25_index = BM25Okapi(bm25_corpus)
                    print(
                        f"✓ Loaded BM25 index ({len(bm25_corpus)} documents)",
                        flush=True,
                    )
                    return

        # Rebuild from database
        rebuild_bm25_index()
    except Exception as e:
        print(f"⚠️ Could not load BM25 index: {e}", flush=True)
        bm25_index = None


def save_bm25_index():
    """Save BM25 index to disk"""
    try:
        data = {
            "corpus": bm25_corpus,
            "doc_ids": bm25_doc_ids,
            "timestamp": datetime.now().isoformat(),
        }
        with open(BM25_INDEX_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f)
    except Exception as e:
        print(f"⚠️ Could not save BM25 index: {e}", flush=True)


def rebuild_bm25_index():
    """Rebuild BM25 index from database"""
    global bm25_index, bm25_corpus, bm25_doc_ids

    try:
        df = documents_table.to_pandas()
        if len(df) == 0:
            bm25_index = None
            bm25_corpus = []
            bm25_doc_ids = []
            return

        # Group by document and create keyword corpus
        doc_keywords = defaultdict(list)
        for _, row in df.iterrows():
            doc_id = row["document_id"]
            chunk_text = row["text"]
            # Extract keywords (simple approach: tokenize and get unique words)
            tokens = tokenize_text(chunk_text)
            doc_keywords[doc_id].extend(tokens)

        bm25_corpus = []
        bm25_doc_ids = []

        for doc_id, tokens in doc_keywords.items():
            unique_tokens = list(set(tokens))
            bm25_corpus.append(unique_tokens)
            bm25_doc_ids.append(doc_id)

        if bm25_corpus:
            bm25_index = BM25Okapi(bm25_corpus)
            save_bm25_index()
            print(f"✓ Rebuilt BM25 index ({len(bm25_corpus)} documents)", flush=True)
        else:
            bm25_index = None
    except Exception as e:
        print(f"⚠️ Could not rebuild BM25 index: {e}", flush=True)
        bm25_index = None


def add_to_bm25_index(doc_id: str, chunks: List[dict]):
    """Add document to BM25 index"""
    global bm25_index

    try:
        # Combine all chunk texts
        all_tokens = []
        for chunk in chunks:
            tokens = tokenize_text(chunk["text"])
            all_tokens.extend(tokens)

        # Remove duplicates for this document
        unique_tokens = list(set(all_tokens))

        # Check if document already exists in index
        if doc_id in bm25_doc_ids:
            idx = bm25_doc_ids.index(doc_id)
            bm25_corpus[idx] = unique_tokens
        else:
            bm25_doc_ids.append(doc_id)
            bm25_corpus.append(unique_tokens)

        # Rebuild BM25 index
        if bm25_corpus:
            bm25_index = BM25Okapi(bm25_corpus)
            save_bm25_index()
    except Exception as e:
        print(f"⚠️ Could not add to BM25 index: {e}", flush=True)


def extract_keywords(text: str, max_keywords: int = 20) -> str:
    """Extract keywords from text for BM25 indexing"""
    tokens = tokenize_text(text)
    # Simple keyword extraction: most frequent words (excluding common stop words)
    stop_words = {
        "the",
        "a",
        "an",
        "and",
        "or",
        "but",
        "in",
        "on",
        "at",
        "to",
        "for",
        "of",
        "with",
        "by",
        "is",
        "are",
        "was",
        "were",
        "be",
        "been",
        "have",
        "has",
        "had",
        "do",
        "does",
        "did",
        "will",
        "would",
        "could",
        "should",
        "may",
        "might",
        "must",
        "can",
        "this",
        "that",
        "these",
        "those",
    }

    word_freq = {}
    for token in tokens:
        if len(token) > 2 and token not in stop_words:
            word_freq[token] = word_freq.get(token, 0) + 1

    # Get top keywords
    sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
    keywords = [word for word, freq in sorted_words[:max_keywords]]

    return " ".join(keywords)


def extract_text_from_pdf(file_path: str) -> List[dict]:
    """Extract text from PDF with page numbers"""
    pages = []

    if UNSTRUCTURED_AVAILABLE:
        try:
            elements = partition_pdf(
                filename=file_path, strategy="fast", include_page_breaks=True
            )
            current_page = 1
            page_texts = {}

            for element in elements:
                page_num = element.metadata.page_number or current_page
                if page_num not in page_texts:
                    page_texts[page_num] = []
                page_texts[page_num].append(str(element))

            for page_num in sorted(page_texts.keys()):
                pages.append(
                    {"page": page_num, "text": "\n".join(page_texts[page_num])}
                )
            return pages
        except Exception as e:
            print(f"Unstructured PDF failed: {e}, falling back", flush=True)

    if PYPDF2_AVAILABLE:
        try:
            with open(file_path, "rb") as f:
                pdf_reader = PyPDF2.PdfReader(f)
                total_pages = len(pdf_reader.pages)
                for page_num in range(total_pages):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    pages.append({"page": page_num + 1, "text": text})
            return pages
        except Exception as e:
            raise Exception(f"PDF extraction failed: {e}")
    else:
        raise Exception("No PDF extraction library available")


def extract_text_from_docx(file_path: str) -> List[dict]:
    """Extract text from Word document"""
    pages = []

    if UNSTRUCTURED_AVAILABLE:
        try:
            elements = partition_docx(
                filename=file_path, infer_table_structure=True, include_page_breaks=True
            )
            current_page = 1
            page_texts = {}

            for element in elements:
                if hasattr(element.metadata, "page_number"):
                    page_num = element.metadata.page_number
                else:
                    page_num = current_page
                if page_num not in page_texts:
                    page_texts[page_num] = []
                page_texts[page_num].append(str(element))

            for page_num in sorted(page_texts.keys()):
                pages.append(
                    {"page": page_num, "text": "\n".join(page_texts[page_num])}
                )
            return pages
        except Exception as e:
            print(f"Unstructured DOCX failed: {e}, falling back", flush=True)

    if PYTHON_DOCX_AVAILABLE:
        try:
            doc = docx.Document(file_path)
            full_text = [para.text for para in doc.paragraphs]
            pages.append({"page": 1, "text": "\n".join(full_text)})
            return pages
        except Exception as e:
            raise Exception(f"DOCX extraction failed: {e}")
    else:
        raise Exception("No DOCX extraction library available")


def extract_text_from_excel(file_path: str) -> List[dict]:
    """Extract text from Excel spreadsheet"""
    pages = []

    if PANDAS_AVAILABLE:
        try:
            # Read all sheets
            xl_file = pd.ExcelFile(file_path)
            full_text = []

            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(xl_file, sheet_name=sheet_name)
                full_text.append(f"Sheet: {sheet_name}")
                full_text.append(df.to_string(index=False))
                full_text.append("\n")

            pages.append({"page": 1, "text": "\n".join(full_text)})
            return pages
        except Exception as e:
            raise Exception(f"Excel extraction failed: {e}")
    elif OPENPYXL_AVAILABLE:
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            full_text = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                full_text.append(f"Sheet: {sheet_name}")
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        full_text.append(row_text)
                full_text.append("\n")

            pages.append({"page": 1, "text": "\n".join(full_text)})
            return pages
        except Exception as e:
            raise Exception(f"Excel extraction failed: {e}")
    else:
        raise Exception(
            "No Excel extraction library available. Install: pip install pandas openpyxl"
        )


def extract_text_from_csv(file_path: str) -> List[dict]:
    """Extract text from CSV file"""
    pages = []

    if PANDAS_AVAILABLE:
        try:
            df = pd.read_csv(file_path)
            pages.append({"page": 1, "text": df.to_string(index=False)})
            return pages
        except Exception as e:
            raise Exception(f"CSV extraction failed: {e}")
    else:
        # Fallback to standard library
        try:
            import csv

            full_text = []
            with open(file_path, "r", encoding="utf-8") as f:
                reader = csv.reader(f)
                for row in reader:
                    full_text.append(" | ".join(row))
            pages.append({"page": 1, "text": "\n".join(full_text)})
            return pages
        except Exception as e:
            raise Exception(f"CSV extraction failed: {e}")


def extract_text_from_powerpoint(file_path: str) -> List[dict]:
    """Extract text from PowerPoint presentation"""
    pages = []

    if PYTHON_PPTX_AVAILABLE:
        try:
            prs = Presentation(file_path)

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                    # Extract from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(cell.text for cell in row.cells)
                            slide_text.append(row_text)

                if slide_text:
                    pages.append(
                        {
                            "page": slide_num,
                            "text": f"Slide {slide_num}:\n" + "\n".join(slide_text),
                        }
                    )

            return pages
        except Exception as e:
            raise Exception(f"PowerPoint extraction failed: {e}")
    else:
        raise Exception(
            "No PowerPoint extraction library available. Install: pip install python-pptx"
        )


def extract_text_from_text(file_path: str, encoding: str = "utf-8") -> List[dict]:
    """Extract text from plain text files"""
    pages = []

    try:
        with open(file_path, "r", encoding=encoding) as f:
            text = f.read()
        pages.append({"page": 1, "text": text})
        return pages
    except UnicodeDecodeError:
        # Try with different encoding
        with open(file_path, "r", encoding="latin-1") as f:
            text = f.read()
        pages.append({"page": 1, "text": text})
        return pages


def extract_text_auto(file_path: str) -> List[dict]:
    """Auto-detect file type and extract text"""
    file_ext = Path(file_path).suffix.lower()

    if UNSTRUCTURED_AVAILABLE and file_ext in [
        ".pdf",
        ".docx",
        ".doc",
        ".html",
        ".htm",
        ".xml",
        ".rtf",
        ".epub",
        ".odt",
    ]:
        try:
            elements = partition(filename=file_path)
            return [{"page": 1, "text": "\n".join(str(el) for el in elements)}]
        except Exception as e:
            print(
                f"Unstructured auto-partition failed: {e}, trying specific parser",
                flush=True,
            )

    # Fallback to specific parsers
    if file_ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif file_ext in [".docx", ".doc"]:
        return extract_text_from_docx(file_path)
    elif file_ext in [".xlsx", ".xls"]:
        return extract_text_from_excel(file_path)
    elif file_ext == ".csv":
        return extract_text_from_csv(file_path)
    elif file_ext in [".pptx", ".ppt"]:
        return extract_text_from_powerpoint(file_path)
    elif file_ext in [
        ".txt",
        ".md",
        ".markdown",
        ".rst",
        ".log",
        ".tex",
        ".css",
        ".js",
        ".py",
        ".java",
        ".cpp",
        ".c",
        ".h",
        ".sql",
        ".sh",
        ".bat",
        ".ps1",
        ".yaml",
        ".yml",
        ".json",
        ".xml",
        ".html",
        ".htm",
    ]:
        return extract_text_from_text(file_path)
    else:
        # Try as text file
        return extract_text_from_text(file_path)


def chunk_pages(pages: List[dict]) -> List[dict]:
    """Split pages into semantic chunks"""
    chunks = []
    chunk_index = 0

    for page in pages:
        page_text = page["text"]
        page_num = page["page"]

        if not page_text.strip():
            continue

        page_chunks = text_splitter.split_text(page_text)

        for chunk_text in page_chunks:
            chunks.append(
                {
                    "chunk_index": chunk_index,
                    "page_number": page_num,
                    "text": chunk_text,
                }
            )
            chunk_index += 1

    return chunks


def reciprocal_rank_fusion(
    results_lists: List[List[Tuple[str, float]]], k: int = 60
) -> List[Tuple[str, float]]:
    """
    Combine multiple ranked lists using Reciprocal Rank Fusion

    Args:
        results_lists: List of ranked lists, each containing (doc_id, score) tuples
        k: RRF constant (default 60)

    Returns:
        List of (doc_id, rrf_score) tuples sorted by score
    """
    rrf_scores = defaultdict(float)

    for results in results_lists:
        for rank, (doc_id, _) in enumerate(results, start=1):
            rrf_scores[doc_id] += 1.0 / (k + rank)

    # Sort by RRF score
    sorted_results = sorted(rrf_scores.items(), key=lambda x: x[1], reverse=True)
    return sorted_results


@server.list_tools()
async def list_tools() -> list[Tool]:
    """List available document RAG tools"""
    return [
        Tool(
            name="upload_document",
            description="Upload and index a document (PDF, DOCX, XLSX, PPTX, TXT, MD, CSV, and more) for semantic search",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Absolute path to the document file",
                    },
                    "document_name": {
                        "type": "string",
                        "description": "Display name for the document (optional, defaults to filename)",
                    },
                    "metadata": {
                        "type": "object",
                        "description": "Optional metadata (author, category, tags, etc.)",
                        "default": {},
                    },
                },
                "required": ["file_path"],
            },
        ),
        Tool(
            name="search_documents",
            description="Search indexed documents using hybrid search (BM25 + semantic similarity) with source citations",
            inputSchema={
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "Natural language query to search for",
                    },
                    "top_k": {
                        "type": "integer",
                        "description": "Number of results to return",
                        "default": 5,
                        "minimum": 1,
                        "maximum": 20,
                    },
                    "document_filter": {
                        "type": "string",
                        "description": "Optional: filter by specific document ID",
                    },
                    "search_type": {
                        "type": "string",
                        "enum": ["hybrid", "semantic", "keyword"],
                        "description": "Search type: hybrid (BM25+vectors), semantic (vectors only), or keyword (BM25 only)",
                        "default": "hybrid",
                    },
                    "min_score": {
                        "type": "number",
                        "description": "Minimum similarity score (0-1)",
                        "default": 0.2,
                        "minimum": 0,
                        "maximum": 1,
                    },
                },
                "required": ["query"],
            },
        ),
        Tool(
            name="list_documents",
            description="List all indexed documents with metadata and versions",
            inputSchema={
                "type": "object",
                "properties": {
                    "limit": {
                        "type": "integer",
                        "description": "Maximum documents to list",
                        "default": 50,
                        "minimum": 1,
                        "maximum": 100,
                    }
                },
            },
        ),
        Tool(
            name="delete_document",
            description="Delete a document and all its versions from the index",
            inputSchema={
                "type": "object",
                "properties": {
                    "document_id": {
                        "type": "string",
                        "description": "ID of the document to delete",
                    }
                },
                "required": ["document_id"],
            },
        ),
        Tool(
            name="get_document_versions",
            description="Get version history for a specific document",
            inputSchema={
                "type": "object",
                "properties": {
                    "document_id": {
                        "type": "string",
                        "description": "ID of the document",
                    }
                },
                "required": ["document_id"],
            },
        ),
        Tool(
            name="get_document_stats",
            description="Get statistics about indexed documents",
            inputSchema={"type": "object", "properties": {}},
        ),
        Tool(
            name="chat_with_documents",
            description="Retrieve relevant document chunks for a query (to be used with Claude's context)",
            inputSchema={
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "User question about the documents",
                    },
                    "document_ids": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Optional: limit search to specific documents",
                        "default": [],
                    },
                    "top_k": {
                        "type": "integer",
                        "description": "Number of chunks to retrieve",
                        "default": 5,
                        "minimum": 1,
                        "maximum": 10,
                    },
                },
                "required": ["query"],
            },
        ),
    ]


@server.call_tool()
async def call_tool(name: str, arguments: Any) -> list[TextContent]:
    """Handle tool calls"""
    global documents_table, versions_table, db

    try:
        if name == "upload_document":
            file_path = arguments["file_path"]
            document_name = arguments.get("document_name", Path(file_path).name)
            metadata = arguments.get("metadata", {})

            # Validate file exists
            if not os.path.exists(file_path):
                return [
                    TextContent(type="text", text=f"❌ File not found: {file_path}")
                ]

            # Determine file type
            file_ext = Path(file_path).suffix.lower()

            if file_ext not in SUPPORTED_FORMATS:
                return [
                    TextContent(
                        type="text",
                        text=f"❌ Unsupported file type: {file_ext}. Supported: {', '.join(SUPPORTED_FORMATS.keys())}",
                    )
                ]

            # Generate document ID
            with open(file_path, "rb") as f:
                file_hash = hashlib.md5(f.read()).hexdigest()[:16]
            document_id = f"doc_{file_hash}"

            # Check for existing versions
            df_versions = versions_table.to_pandas()
            existing_versions = df_versions[df_versions["document_id"] == document_id]

            if not existing_versions.empty:
                # Increment version
                version_number = existing_versions["version_number"].max() + 1
            else:
                version_number = 1

            # Extract text based on file type
            print(f"Processing {file_ext} file: {file_path}", flush=True)

            try:
                pages = extract_text_auto(file_path)
            except Exception as e:
                return [
                    TextContent(
                        type="text", text=f"❌ Failed to extract text: {str(e)}"
                    )
                ]

            total_pages = len(pages)
            print(f"Extracted {total_pages} pages", flush=True)

            # Chunk the text
            chunks = chunk_pages(pages)
            print(f"Created {len(chunks)} chunks", flush=True)

            # Generate embeddings and store
            records = []
            chunk_texts = [c["text"] for c in chunks]

            with torch.no_grad():
                embeddings = embed_model.encode(
                    chunk_texts,
                    batch_size=32,
                    convert_to_tensor=False,
                    show_progress_bar=False,
                )

            # Extract keywords for BM25
            all_text = " ".join(chunk_texts)
            keywords = extract_keywords(all_text)

            for i, chunk in enumerate(chunks):
                record = {
                    "id": f"{document_id}_v{version_number}_chunk_{chunk['chunk_index']}",
                    "text": chunk["text"],
                    "vector": embeddings[i].tolist(),
                    "document_id": document_id,
                    "document_name": document_name,
                    "chunk_index": chunk["chunk_index"],
                    "page_number": chunk["page_number"],
                    "total_pages": total_pages,
                    "timestamp": datetime.now().isoformat(),
                    "file_type": SUPPORTED_FORMATS[file_ext],
                    "version": version_number,
                    "keywords": keywords,
                    "metadata": json.dumps(metadata),
                }
                records.append(record)

            # Store in LanceDB
            documents_table.add(records)

            # Add version record
            version_record = {
                "version_id": f"{document_id}_v{version_number}",
                "document_id": document_id,
                "document_name": document_name,
                "version_number": version_number,
                "timestamp": datetime.now().isoformat(),
                "file_hash": file_hash,
                "file_path": file_path,
                "chunk_count": len(chunks),
                "metadata": json.dumps(metadata),
            }
            versions_table.add([version_record])

            # Update BM25 index
            add_to_bm25_index(document_id, chunks)

            return [
                TextContent(
                    type="text",
                    text=f"✓ Indexed '{document_name}' (Version {version_number})\n"
                    f"  • Document ID: {document_id}\n"
                    f"  • Pages: {total_pages}\n"
                    f"  • Chunks: {len(chunks)}\n"
                    f"  • File type: {SUPPORTED_FORMATS[file_ext]}\n"
                    f"  • Storage: {DOCUMENTS_DIR}",
                )
            ]

        elif name == "search_documents":
            query = arguments["query"]
            top_k = arguments.get("top_k", 5)
            document_filter = arguments.get("document_filter")
            search_type = arguments.get("search_type", "hybrid")
            min_score = arguments.get("min_score", 0.2)

            results = []

            # Semantic search
            if search_type in ["hybrid", "semantic"]:
                with torch.no_grad():
                    query_embedding = embed_model.encode(
                        [query], convert_to_tensor=False, show_progress_bar=False
                    )[0].tolist()

                semantic_results = (
                    documents_table.search(query_embedding).limit(top_k * 2).to_list()
                )
                semantic_ranked = []

                seen_docs = set()
                for r in semantic_results:
                    score = 1 / (1 + r.get("_distance", 1))
                    if score < min_score:
                        continue
                    if document_filter and r["document_id"] != document_filter:
                        continue

                    doc_id = r["document_id"]
                    if doc_id not in seen_docs:
                        semantic_ranked.append((doc_id, score))
                        seen_docs.add(doc_id)

                    if len(semantic_ranked) >= top_k * 2:
                        break

                results.append(semantic_ranked)

            # Keyword search (BM25)
            if (
                search_type in ["hybrid", "keyword"]
                and bm25_index
                and RANK_BM25_AVAILABLE
            ):
                query_tokens = tokenize_text(query)
                bm25_scores = bm25_index.get_scores(query_tokens)

                # Get top results
                bm25_results = []
                for idx, score in enumerate(bm25_scores):
                    if score > 0:
                        doc_id = bm25_doc_ids[idx]
                        if document_filter and doc_id != document_filter:
                            continue
                        # Normalize score
                        normalized_score = min(score / 10.0, 1.0)  # Rough normalization
                        bm25_results.append((doc_id, normalized_score))

                bm25_results.sort(key=lambda x: x[1], reverse=True)
                results.append(bm25_results[: top_k * 2])

            # Fuse results using RRF if hybrid
            if search_type == "hybrid" and len(results) >= 2:
                fused_results = reciprocal_rank_fusion(results, k=60)
                final_results = fused_results[:top_k]
            elif len(results) == 1:
                final_results = results[0][:top_k]
            else:
                final_results = []

            if not final_results:
                return [TextContent(type="text", text="No relevant documents found.")]

            # Get chunk details for top results
            df = documents_table.to_pandas()
            result_text = f"Found {len(final_results)} relevant documents:\n\n"

            for i, (doc_id, score) in enumerate(final_results, 1):
                doc_chunks = df[df["document_id"] == doc_id]
                if len(doc_chunks) == 0:
                    continue

                doc_name = doc_chunks.iloc[0]["document_name"]
                version = doc_chunks.iloc[0]["version"]
                file_type = doc_chunks.iloc[0]["file_type"]

                result_text += f"{i}. 📄 {doc_name} (v{version}, {file_type})\n"
                result_text += f"   Score: {score:.3f} | ID: {doc_id}\n"

                # Show top chunk preview
                top_chunk = doc_chunks.iloc[0]
                preview = (
                    top_chunk["text"][:200] + "..."
                    if len(top_chunk["text"]) > 200
                    else top_chunk["text"]
                )
                result_text += f"   Preview: {preview}\n\n"

            return [TextContent(type="text", text=result_text)]

        elif name == "list_documents":
            limit = arguments.get("limit", 50)

            df = versions_table.to_pandas()

            if len(df) == 0:
                return [TextContent(type="text", text="No documents indexed yet.")]

            # Get latest version of each document
            latest_versions = (
                df.sort_values("version_number", ascending=False)
                .groupby("document_id")
                .first()
                .reset_index()
            )
            latest_versions = latest_versions.head(limit)

            result_text = f"📚 Indexed Documents ({len(latest_versions)} unique):\n\n"
            for _, doc in latest_versions.iterrows():
                # Count versions
                version_count = len(df[df["document_id"] == doc["document_id"]])

                result_text += f"• {doc['document_name']}\n"
                result_text += f"  ID: {doc['document_id']}\n"
                result_text += f"  Latest: v{doc['version_number']} | Total versions: {version_count}\n"
                result_text += f"  Chunks: {doc['chunk_count']} | Indexed: {doc['timestamp'][:10]}\n\n"

            return [TextContent(type="text", text=result_text)]

        elif name == "delete_document":
            document_id = arguments["document_id"]

            df_docs = documents_table.to_pandas()
            df_versions = versions_table.to_pandas()

            if (
                document_id not in df_docs["document_id"].values
                and document_id not in df_versions["document_id"].values
            ):
                return [
                    TextContent(
                        type="text", text=f"❌ Document '{document_id}' not found."
                    )
                ]

            # Filter out the document from documents table
            df_docs_filtered = df_docs[df_docs["document_id"] != document_id]

            # Recreate documents table
            db.drop_table("documents")

            schema = pa.schema(
                [
                    pa.field("id", pa.string()),
                    pa.field("text", pa.string()),
                    pa.field("vector", pa.list_(pa.float32(), 384)),
                    pa.field("document_id", pa.string()),
                    pa.field("document_name", pa.string()),
                    pa.field("chunk_index", pa.int32()),
                    pa.field("page_number", pa.int32()),
                    pa.field("total_pages", pa.int32()),
                    pa.field("timestamp", pa.string()),
                    pa.field("file_type", pa.string()),
                    pa.field("version", pa.int32()),
                    pa.field("keywords", pa.string()),
                    pa.field("metadata", pa.string()),
                ]
            )

            documents_table = db.create_table("documents", schema=schema)

            if len(df_docs_filtered) > 0:
                records = df_docs_filtered.to_dict("records")
                documents_table.add(records)

            # Filter out the document from versions table
            df_versions_filtered = df_versions[
                df_versions["document_id"] != document_id
            ]

            db.drop_table("document_versions")

            version_schema = pa.schema(
                [
                    pa.field("version_id", pa.string()),
                    pa.field("document_id", pa.string()),
                    pa.field("document_name", pa.string()),
                    pa.field("version_number", pa.int32()),
                    pa.field("timestamp", pa.string()),
                    pa.field("file_hash", pa.string()),
                    pa.field("file_path", pa.string()),
                    pa.field("chunk_count", pa.int32()),
                    pa.field("metadata", pa.string()),
                ]
            )

            versions_table = db.create_table("document_versions", schema=version_schema)

            if len(df_versions_filtered) > 0:
                records = df_versions_filtered.to_dict("records")
                versions_table.add(records)

            # Remove from BM25 index
            if document_id in bm25_doc_ids:
                idx = bm25_doc_ids.index(document_id)
                bm25_doc_ids.pop(idx)
                bm25_corpus.pop(idx)
                if bm25_corpus:
                    bm25_index = BM25Okapi(bm25_corpus)
                else:
                    bm25_index = None
                save_bm25_index()

            return [
                TextContent(
                    type="text",
                    text=f"✓ Deleted document '{document_id}' and all its versions.",
                )
            ]

        elif name == "get_document_versions":
            document_id = arguments["document_id"]

            df = versions_table.to_pandas()
            doc_versions = df[df["document_id"] == document_id].sort_values(
                "version_number", ascending=False
            )

            if len(doc_versions) == 0:
                return [
                    TextContent(
                        type="text",
                        text=f"❌ No versions found for document '{document_id}'.",
                    )
                ]

            result_text = (
                f"📋 Version History for {doc_versions.iloc[0]['document_name']}:\n\n"
            )
            for _, version in doc_versions.iterrows():
                result_text += f"Version {version['version_number']}:\n"
                result_text += f"  • Date: {version['timestamp'][:19]}\n"
                result_text += f"  • Chunks: {version['chunk_count']}\n"
                result_text += f"  • Hash: {version['file_hash']}\n\n"

            return [TextContent(type="text", text=result_text)]

        elif name == "get_document_stats":
            df_docs = documents_table.to_pandas()
            df_versions = versions_table.to_pandas()

            if len(df_docs) == 0:
                return [TextContent(type="text", text="No documents indexed yet.")]

            stats = {
                "Total chunks": len(df_docs),
                "Unique documents": df_docs["document_id"].nunique(),
                "Total versions": len(df_versions),
                "File types": df_docs["file_type"].value_counts().to_dict(),
                "Average chunk length": int(df_docs["text"].str.len().mean()),
                "Total text (chars)": df_docs["text"].str.len().sum(),
                "BM25 indexed docs": len(bm25_doc_ids) if bm25_doc_ids else 0,
            }

            result_text = "📊 Document Statistics:\n\n"
            for key, value in stats.items():
                result_text += f"{key}: {value}\n"

            return [TextContent(type="text", text=result_text)]

        elif name == "chat_with_documents":
            query = arguments["query"]
            document_ids = arguments.get("document_ids", [])
            top_k = arguments.get("top_k", 5)

            with torch.no_grad():
                query_embedding = embed_model.encode(
                    [query], convert_to_tensor=False, show_progress_bar=False
                )[0].tolist()

            results = documents_table.search(query_embedding).limit(top_k * 2).to_list()

            filtered = []
            for r in results:
                if document_ids and r["document_id"] not in document_ids:
                    continue

                score = 1 / (1 + r.get("_distance", 1))
                filtered.append(
                    {
                        "text": r["text"],
                        "document_name": r["document_name"],
                        "page_number": r["page_number"],
                        "score": round(score, 3),
                        "version": r["version"],
                    }
                )

                if len(filtered) >= top_k:
                    break

            if not filtered:
                return [
                    TextContent(
                        type="text", text="No relevant context found in documents."
                    )
                ]

            context_text = f"Relevant document passages for: '{query}'\n\n"
            for i, doc in enumerate(filtered, 1):
                context_text += f"[{i}] From '{doc['document_name']}' (v{doc['version']}, Page {doc['page_number']}):\n"
                context_text += f"{doc['text']}\n\n"

            return [TextContent(type="text", text=context_text)]

        else:
            return [TextContent(type="text", text=f"Unknown tool: {name}")]

    except Exception as e:
        print(f"✗ Tool execution error: {e}", flush=True)
        import traceback

        traceback.print_exc(file=sys.stderr)
        return [TextContent(type="text", text=f"Error: {str(e)}")]


async def main():
    """Main entry point"""
    print("Starting Document RAG MCP Server...", flush=True)
    print(f"Python version: {sys.version}", flush=True)
    print(f"Platform: {platform.system()} {platform.release()}", flush=True)

    try:
        init_database()
        print("✓ Initialization complete", flush=True)

        async with stdio_server() as (read_stream, write_stream):
            print("✓ Document RAG Server running and ready", flush=True)
            await server.run(
                read_stream, write_stream, server.create_initialization_options()
            )
    except Exception as e:
        print(f"✗ Fatal error: {e}", flush=True)
        import traceback

        traceback.print_exc(file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        print("\n✓ Server stopped by user", flush=True)
    except Exception as e:
        print(f"✗ Startup error: {e}", flush=True)
        import traceback

        traceback.print_exc(file=sys.stderr)
        sys.exit(1)
